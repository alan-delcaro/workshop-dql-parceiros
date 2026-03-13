#!/usr/bin/env python3
"""
Gerador de PPTX editável — Workshop DQL para Parceiros Dynatrace.
Todos os elementos (texto, tabelas, código) são nativos do PowerPoint.

Uso:
    python3 slides/gerar_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── Paleta Dynatrace ─────────────────────────────────────────────────────────
TEAL       = RGBColor(0x00, 0xB4, 0xD8)
DARK_TEAL  = RGBColor(0x00, 0x77, 0xB6)
NAVY       = RGBColor(0x02, 0x3E, 0x8A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
DARK_TEXT  = RGBColor(0x16, 0x1B, 0x22)
GRAY_TEXT  = RGBColor(0x57, 0x60, 0x6E)
RED        = RGBColor(0xCF, 0x22, 0x2E)
GREEN      = RGBColor(0x1A, 0x7F, 0x37)


# ─── Helpers ──────────────────────────────────────────────────────────────────

def new_blank(prs):
    """Slide em branco sem placeholders."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)
    return slide


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def txbox(slide, text, l, t, w, h,
          size=18, bold=False, color=DARK_TEXT,
          align=PP_ALIGN.LEFT, font="Calibri", wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.italic = italic
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = font
    return tb


def code_block(slide, code, l=0.4, t=1.3, w=12.5, h=4.0, size=10.5):
    """Bloco de código: retângulo cinza + texto monospace editável."""
    rect = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = LIGHT_GRAY
    rect.line.fill.background()

    tb = slide.shapes.add_textbox(
        Inches(l + 0.12), Inches(t + 0.1),
        Inches(w - 0.24), Inches(h - 0.2)
    )
    tf = tb.text_frame
    tf.word_wrap = False
    lines = code.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = "Courier New"
        run.font.size = Pt(size)
        run.font.color.rgb = DARK_TEXT
    return tb


def table(slide, headers, rows, l=0.4, t=1.3, w=12.5, col_widths=None,
          row_h=0.42, hdr_size=13, cell_size=12):
    """Tabela nativa do PowerPoint com cabeçalho teal."""
    nr = len(rows) + 1
    nc = len(headers)
    total_h = row_h * nr
    gf = slide.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(w), Inches(total_h))
    tbl = gf.table

    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(cw)

    def fill_cell(cell, text, bold=False, bg=None, color=DARK_TEXT, size=13, align=PP_ALIGN.LEFT):
        cell.text = ""
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.bold = bold
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

    for i, h in enumerate(headers):
        fill_cell(tbl.cell(0, i), h, bold=True, bg=DARK_TEAL, color=WHITE,
                  size=hdr_size, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        bg = LIGHT_GRAY if r % 2 == 0 else WHITE
        for c, val in enumerate(row):
            fill_cell(tbl.cell(r + 1, c), str(val), bg=bg, size=cell_size)
    return tbl


def divider_line(slide):
    """Barra vertical teal à esquerda e linha abaixo do título."""
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.08), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    ln = slide.shapes.add_shape(1, Inches(0.25), Inches(0.98), Inches(12.83), Inches(0.04))
    ln.fill.solid()
    ln.fill.fore_color.rgb = TEAL
    ln.line.fill.background()


# ─── Templates de slides ──────────────────────────────────────────────────────

def slide_cover(prs, title, subtitle=None, detail=None):
    s = new_blank(prs)
    set_bg(s, NAVY)
    txbox(s, title, 0.5, 2.4, 12.3, 1.1,
          size=38, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    if subtitle:
        txbox(s, subtitle, 0.5, 3.55, 12.3, 0.7,
              size=22, color=WHITE, align=PP_ALIGN.CENTER)
    if detail:
        txbox(s, detail, 0.5, 4.35, 12.3, 0.5,
              size=14, color=GRAY_TEXT, align=PP_ALIGN.CENTER)
    return s


def slide_module(prs, num, title, subtitle=None):
    s = new_blank(prs)
    set_bg(s, DARK_TEAL)
    if num:
        txbox(s, num, 0.5, 1.9, 12.3, 0.6,
              size=18, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s, title, 0.5, 2.5, 12.3, 1.1,
          size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        txbox(s, subtitle, 0.5, 3.65, 12.3, 0.55,
              size=15, color=WHITE, align=PP_ALIGN.CENTER)
    return s


def slide_content(prs, title):
    s = new_blank(prs)
    set_bg(s, WHITE)
    divider_line(s)
    txbox(s, title, 0.25, 0.15, 12.83, 0.75,
          size=24, bold=True, color=DARK_TEAL)
    return s


def note(s, text, t=5.5, color=GRAY_TEXT, size=13, bold=False):
    txbox(s, text, 0.25, t, 12.83, 0.7, size=size, color=color, bold=bold)


# ─── Criar apresentação ───────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


# ── Slide 1: Título ───────────────────────────────────────────────────────────
slide_cover(prs,
    "DQL para Parceiros Dynatrace",
    "Dynatrace Query Language — Do Básico ao Avançado",
    "Duração: 1h45–2h  |  Formato: Demo ao vivo + Labs para depois")

# ── Slide 2: Agenda ───────────────────────────────────────────────────────────
s = slide_content(prs, "Agenda")
table(s,
    headers=["#", "Módulo", "Duração"],
    rows=[
        ["0", "Introdução — Grail e DQL", "10 min"],
        ["1", "DQL 100 — Fundamentos", "25 min"],
        ["2", "DQL 200 — Agregação e Parsing", "25 min"],
        ["3", "DQL 300 — Enriquecimento e Joins", "25 min"],
        ["4", "Casos de Uso Práticos", "15 min"],
        ["—", "Q&A e próximos passos", "5 min"],
    ],
    l=1.8, t=1.15, w=9.8, col_widths=[0.5, 7.8, 1.5])
note(s, "Labs hands-on são entregues como material complementar — executados no ritmo do participante.")


# ─── MÓDULO 0 ────────────────────────────────────────────────────────────────

# Slide 3
slide_module(prs, "Módulo 0", "O que é Grail e por que DQL?")

# Slide 4
s = slide_content(prs, "O Grail — Data Lakehouse Nativo")
code_block(s, """\
Logs      ──►  ┌──────────────────────────────┐
Métricas  ──►  │            GRAIL             │
Spans     ──►  │   sem pré-agregação          │
BizEvents ──►  │   sem schema fixo            │
Entidades ──►  │   auto-correlacionado        │
               └──────────────┬───────────────┘
                               │  queryable via DQL
                    ┌──────────▼──────────┐
                    │  Notebook / Alert   │
                    │  Dashboard / API    │
                    └─────────────────────┘""",
    t=1.12, h=3.8)
note(s, "Sem index tuning.  Sem partitions.  Sem pré-computação.", t=5.1, color=DARK_TEAL, bold=True, size=16)

# Slide 5
s = slide_content(prs, "Por que DQL Unifica Tudo?")
table(s,
    headers=["Antes", "Com DQL"],
    rows=[
        ["Log Viewer (busca limitada)", "fetch logs | filter ..."],
        ["Metrics Explorer (só clique)", "timeseries avg(cpu)"],
        ["USQL (só sessões)", "fetch bizevents | summarize ..."],
        ["Cruzar logs + traces = impossível", "fetch spans | join [fetch logs]"],
    ],
    l=0.25, t=1.12, w=12.83, col_widths=[6.4, 6.43])
note(s, "Uma linguagem.  Todas as fontes.  Um resultado.",
    color=DARK_TEAL, bold=True, size=18)

# Slide 6
s = slide_content(prs, "Fontes de Dados (Record Sources)")
table(s,
    headers=["Record Source", "O que contém"],
    rows=[
        ["logs", "Registros de log"],
        ["metrics", "Métricas OOTB e custom"],
        ["spans", "Traces distribuídos"],
        ["bizevents", "Eventos de negócio"],
        ["events", "Deploys, mudanças, alertas"],
        ["dt.entity.*", "Modelo de entidades Dynatrace"],
    ],
    l=1.8, t=1.12, w=9.8, col_widths=[3.5, 6.3])

# Slide 7
s = slide_content(prs, "Conceito de Pipeline")
code_block(s, """\
fetch logs                            // 1. busca dados
| filter loglevel == "ERROR"          // 2. filtra
| summarize count = count(),          // 3. agrega
    by: {dt.entity.service}
| sort count, direction: "descending" // 4. ordena
| limit 10                            // 5. limita""",
    t=1.12, h=3.0)
note(s, "Pense como um cano: dados entram brutos, cada | os transforma, o resultado sai refinado.", t=4.35)


# ─── MÓDULO 1 ────────────────────────────────────────────────────────────────

# Slide 8
slide_module(prs, "Módulo 1", "DQL 100 — Fundamentos",
             "fetch  ·  filter  ·  filterOut  ·  timeseries")

# Slide 9: fetch
s = slide_content(prs, "fetch — A Porta de Entrada")
code_block(s, """\
// Sintaxe básica
fetch logs

// Com janela de tempo
fetch logs, from: now() - 2h

// Com limite de scan (boa prática!)
fetch logs, from: now() - 7d, scanLimitGBytes: 1""",
    t=1.12, h=3.3)
note(s, "Unidades de tempo:  s  m  h  d  w", t=4.6, size=14)
note(s, "Sempre defina from: em queries de produção para resultados previsíveis.", t=5.1, size=13)

# Slide 10: filter
s = slide_content(prs, "filter — Mantém o que Passa")
code_block(s, """\
fetch logs, from: now() - 1h
| filter loglevel == "ERROR"           // igual
| filter loglevel in ["ERROR","FATAL"] // lista
| filter duration > 1000               // comparação
| filter isNotNull(dt.entity.service)  // não nulo

// Operadores lógicos:
| filter loglevel == "ERROR" and dt.entity.service != null
| filter loglevel == "ERROR" or loglevel == "WARN"
| filter not(loglevel in ["DEBUG","TRACE"])""",
    t=1.12, h=4.7)

# Slide 11: filterOut
s = slide_content(prs, "filterOut — Remove o que Combina")
code_block(s, """\
fetch logs, from: now() - 1h
| filterOut contains(content, "/health")
    or contains(content, "/ping")
| filterOut loglevel in ["DEBUG","TRACE"]""",
    t=1.12, h=2.8)
note(s, "Use filterOut quando é mais natural descrever o que você não quer.", t=4.2, size=15, color=GRAY_TEXT)

# Slide 12: Funções de Match
s = slide_content(prs, "Funções de Match em Strings")
table(s,
    headers=["Função", "Wildcards", "Uso"],
    rows=[
        ['contains(campo, "val")', "Não", "Substring simples"],
        ['matchesValue(campo, "val*")', "* e ?", "Match com padrão"],
        ['matchesPhrase(campo, "frase")', "Não", "Frase em texto livre"],
        ["startsWith / endsWith", "Não", "Prefixo / sufixo"],
    ],
    l=0.25, t=1.12, w=12.83, col_widths=[5.5, 1.5, 5.83])
code_block(s, """\
| filter matchesValue(dt.entity.service_name, "payment*")
| filter matchesPhrase(content, "connection refused")""",
    t=4.05, h=1.2)

# Slide 13: Controles de Volume
s = slide_content(prs, "Controles de Volume")
code_block(s, """\
// Limita a saída (não o scan!)
| limit 100

// Amostra aleatória — só para exploração
| samplingRatio 0.1         // analisa 10%

// Limita o scan real — usar sempre em janelas grandes
fetch logs, from: now() - 7d, scanLimitGBytes: 1""",
    t=1.12, h=3.7)
note(s, "⚠  samplingRatio retorna estimativas — NUNCA use em alertas ou SLOs.", t=5.0,
    color=RED, bold=True, size=14)

# Slide 14: timeseries
s = slide_content(prs, "timeseries — Métricas no Tempo")
code_block(s, """\
// Média e P95 de CPU dos hosts de produção, a cada 5 min
timeseries {
    avg_cpu = avg(dt.host.cpu.usage),
    p95_cpu = percentile(dt.host.cpu.usage, 95)
},
    filter: {tags contains "production"},
    interval: 5m,
    from: now() - 2h,
    by: {dt.entity.host}""",
    t=1.12, h=3.9)
note(s, "Funções:  avg  ·  sum  ·  min  ·  max  ·  percentile  ·  count", t=5.2, size=14)

# Slide 15: Demo 1
s = slide_content(prs, "Demo 1 — DQL 100 ao Vivo")
code_block(s, """\
// Top 50 erros da última hora — excluindo healthchecks
fetch logs, from: now() - 1h
| filterOut contains(content, "/health") or contains(content, "/ping")
| filter loglevel in ["ERROR", "FATAL"]
| filter isNotNull(dt.entity.service)
| sort timestamp, direction: "descending"
| limit 50""",
    t=1.12, h=3.7)
note(s, "→  Execute no Notebook ao vivo", t=5.05, color=DARK_TEAL, bold=True, size=16)


# ─── MÓDULO 2 ────────────────────────────────────────────────────────────────

# Slide 16
slide_module(prs, "Módulo 2", "DQL 200 — Agregação, Seleção e Parsing",
             "summarize  ·  makeTimeseries  ·  fields  ·  parse")

# Slide 17: summarize
s = slide_content(prs, "summarize — Agregar Dados")
code_block(s, """\
fetch logs, from: now() - 1h
| summarize
    total    = count(),
    errors   = countIf(loglevel == "ERROR"),
    avg_dur  = avg(duration),
    p95_dur  = percentile(duration, 95),
    by: {dt.entity.service}
| sort errors, direction: "descending" """,
    t=1.12, h=3.7)
note(s, "Funções:  count  ·  countIf  ·  countDistinct  ·  avg  ·  sum  ·  min  ·  max  ·  percentile  ·  collectDistinct",
    t=5.05, size=13)

# Slide 18: Error Rate
s = slide_content(prs, "Error Rate por Serviço")
code_block(s, """\
fetch logs, from: now() - 1h
| filter isNotNull(dt.entity.service)
| summarize
    total  = count(),
    errors = countIf(loglevel == "ERROR"),
    by: {dt.entity.service}
| fields dt.entity.service, total, errors,
         rate = round(errors / total * 100, 2)
| sort rate, direction: "descending" """,
    t=1.12, h=4.1)

# Slide 19: makeTimeseries
s = slide_content(prs, "makeTimeseries — Tendência de Qualquer Dado")
code_block(s, """\
fetch logs, from: now() - 2h
| filter loglevel == "ERROR"
| makeTimeseries {
    errors   = countIf(loglevel == "ERROR"),
    warnings = countIf(loglevel == "WARN")
  },
  interval: 5m,
  by: {dt.entity.service}""",
    t=1.12, h=3.5)
table(s,
    headers=["Comando", "Quando usar"],
    rows=[
        ["timeseries", "Apenas métricas nativas"],
        ["makeTimeseries", "Qualquer fonte: logs, spans, bizevents"],
    ],
    l=0.25, t=4.85, w=12.83, col_widths=[3.5, 9.33], row_h=0.38)

# Slide 20: fields
s = slide_content(prs, "fields — Selecionar e Calcular")
code_block(s, """\
fetch logs, from: now() - 1h
| fields
    timestamp,
    service_id = dt.entity.service,
    duration_s = duration / 1000.0,       // cálculo
    is_slow    = duration > 5000,         // booleano
    health     = if(duration > 5000, "SLOW",
                 if(duration > 1000, "OK", "FAST"))  // condicional

// Remover campos:
| fields -trace_id, -span_id

// Explorar campos disponíveis:
| fieldsSummary""",
    t=1.12, h=5.2)

# Slide 21: parse
s = slide_content(prs, "parse — Extrair Campos de Logs")
note(s, 'Log original:  level=ERROR service=payment status=500 duration=2345ms',
    t=1.1, size=13, color=GRAY_TEXT)
code_block(s, """\
fetch logs, from: now() - 1h
| parse content,
    "'level=' WORD:level ' service=' WORD:svc
     ' status=' INT:status ' duration=' LONG:dur_ms 'ms'"
| filter isNotNull(status)
| fields level, svc, status, dur_ms
| sort dur_ms, direction: "descending" """,
    t=1.65, h=3.7)

# Slide 22: Padrões GROK
s = slide_content(prs, "Padrões GROK Mais Usados")
table(s,
    headers=["Padrão", "Captura", "Exemplo"],
    rows=[
        ["WORD", "Sem espaço", "ERROR, payment-api"],
        ["INT", "Inteiro", "200, 500"],
        ["LONG", "Inteiro longo", "timestamps Unix"],
        ["FLOAT", "Decimal", "3.14"],
        ["IPADDR", "Endereço IP", "192.168.1.1"],
        ["LD", "Até próximo padrão", "qualquer texto delimitado"],
        ["DATA", "Qualquer coisa", "texto com espaços"],
        ["TIMESTAMP", "Data/hora ISO", "2024-03-12T10:45:00"],
    ],
    l=0.25, t=1.12, w=12.83, col_widths=[2.8, 3.3, 6.73])

# Slide 23: Demo 2
s = slide_content(prs, "Demo 2 — DQL 200 ao Vivo")
code_block(s, """\
// Análise de performance de endpoints
fetch logs, from: now() - 1h
| filter contains(content, "method=")
| parse content,
    "'method=' WORD:method ' path=' LD:path ' status=' INT:status ' duration=' LONG:ms 'ms'"
| filter isNotNull(status)
| summarize
    requests = count(),
    errors   = countIf(status >= 500),
    p95_ms   = percentile(ms, 95),
    by: {method, path}
| sort p95_ms, direction: "descending"
| limit 15""",
    t=1.12, h=5.3)


# ─── MÓDULO 3 ────────────────────────────────────────────────────────────────

# Slide 24
slide_module(prs, "Módulo 3", "DQL 300 — Enriquecimento e Conexão de Dados",
             "entityAttr  ·  lookup  ·  join  ·  data  ·  append")

# Slide 25: O problema dos Entity IDs
s = slide_content(prs, "O Problema dos Entity IDs")
note(s, 'Logs contêm IDs opacos:  dt.entity.service = "SERVICE-ABCD1234"  ← inútil em dashboard',
    t=1.1, color=RED, size=14)
note(s, 'Queremos:  service_name = "payment-service-prod"  ← legível, utilizável',
    t=1.65, color=GREEN, size=14)
table(s,
    headers=["Abordagem", "Quando usar"],
    rows=[
        ['entityAttr(id, "attr")', "1–2 atributos de entidade, inline, sem subquery"],
        ["lookup [fetch dt.entity.*]", "Múltiplos atributos (nome, estado, tags)"],
        ["join", "Datasets independentes, relação M:N"],
    ],
    l=0.25, t=2.45, w=12.83, col_widths=[4.5, 8.33])

# Slide 26: entityAttr
s = slide_content(prs, "entityAttr — O Mais Simples")
code_block(s, """\
fetch logs, from: now() - 1h
| filter loglevel == "ERROR"
| fields
    timestamp,
    content,
    service = entityAttr(dt.entity.service, "entity.name"),
    host    = entityAttr(dt.entity.host,    "entity.name")
| sort timestamp, direction: "descending"
| limit 30""",
    t=1.12, h=4.1)
note(s, "Rápido. Sem subquery. Ideal para 1–2 atributos.", t=5.35, color=DARK_TEAL, bold=True, size=16)

# Slide 27: lookup
s = slide_content(prs, "lookup — Enriquecimento N:1")
code_block(s, """\
fetch logs, from: now() - 1h
| filter loglevel == "ERROR"
| summarize errors = count(), by: {service_id = dt.entity.service}

| lookup [
    fetch dt.entity.service
    | fields entity.id, entity.name, lifecycleState, tags
  ],
  sourceField: service_id,
  lookupField: entity.id

| fields entity.name, errors, lifecycleState
| sort errors, direction: "descending" """,
    t=1.12, h=5.2)

# Slide 28: data + lookup
s = slide_content(prs, "data + lookup — Tabela de Referência Inline")
code_block(s, """\
fetch logs, from: now() - 2h
| parse content, "LD 'status=' INT:code LD"
| lookup [
    data record(code: 200, label: "Success"),
         record(code: 400, label: "Bad Request"),
         record(code: 401, label: "Unauthorized"),
         record(code: 500, label: "Server Error"),
         record(code: 503, label: "Unavailable")
  ],
  sourceField: code,
  lookupField: code
| summarize count = count(), by: {label}""",
    t=1.12, h=4.6)

# Slide 29: append
s = slide_content(prs, "append — União de Datasets (UNION ALL)")
code_block(s, """\
// Correlacionar erros com deploys na mesma timeline
fetch logs, from: now() - 7d
| filter loglevel == "ERROR"
| fields timestamp, type = "LOG_ERROR", desc = content

| append [
    fetch events, from: now() - 7d
    | filter event.type == "DEPLOYMENT"
    | fields timestamp, type = "DEPLOY", desc = event.name
  ]

| sort timestamp, direction: "descending"
| limit 100""",
    t=1.12, h=4.7)
note(s, "append = UNION ALL    ·    join = JOIN relacional", t=6.0,
    color=DARK_TEAL, bold=True, size=14)

# Slide 30: join
s = slide_content(prs, "join — Cruzar Dois Datasets")
code_block(s, """\
// Error rate combinando dois summarize
fetch logs, from: now() - 1h
| filter loglevel == "ERROR"
| summarize errors = count(), by: {dt.entity.service}

| join [
    fetch logs, from: now() - 1h
    | summarize total = count(), by: {dt.entity.service}
  ],
  on: {dt.entity.service},
  kind: left

| fields
    service   = entityAttr(dt.entity.service, "entity.name"),
    errors, total,
    error_pct = round(errors / total * 100, 2)
| sort error_pct, direction: "descending" """,
    t=1.12, h=5.45)

# Slide 31: Demo 3
s = slide_content(prs, "Demo 3 — Pipeline Mestre DQL 300")
code_block(s, """\
// Service Health Dashboard — última 1 hora
fetch logs, from: now() - 1h
| filter isNotNull(dt.entity.service)
| summarize total = count(), errors = countIf(loglevel == "ERROR"),
    by: {service_id = dt.entity.service}
| lookup [fetch dt.entity.service | fields entity.id, entity.name, lifecycleState],
    sourceField: service_id, lookupField: entity.id
| filter isNotNull(entity.name)
| fields service_name = entity.name, lifecycleState, total, errors,
    error_rate = round(errors / total * 100, 2),
    health = if(errors/total > 0.05, "CRITICAL",
             if(errors/total > 0.01, "DEGRADED", "HEALTHY"))
| sort error_rate, direction: "descending" """,
    t=1.12, h=5.45)


# ─── MÓDULO 4 ────────────────────────────────────────────────────────────────

# Slide 32
slide_module(prs, "Módulo 4", "Casos de Uso Práticos para Parceiros")

# Slide 33: Caso 1
s = slide_content(prs, "Caso 1 — Saúde de Aplicação")
code_block(s, """\
// Top 10 serviços com mais erros — últimas 4h
fetch logs, from: now() - 4h
| filter loglevel in ["ERROR","FATAL"]
| summarize errors = count(), dist_exceptions = countDistinct(exception.type),
    by: {service_id = dt.entity.service}
| lookup [fetch dt.entity.service | fields entity.id, entity.name],
    sourceField: service_id, lookupField: entity.id
| fields entity.name, errors, dist_exceptions
| sort errors, direction: "descending"
| limit 10""",
    t=1.12, h=4.5)

# Slide 34: Caso 2
s = slide_content(prs, "Caso 2 — Latência e SLO Tracking")
code_block(s, """\
// Endpoints violando SLO de 500ms
fetch spans, from: now() - 1h
| filter isNotNull(duration)
| summarize
    requests   = count(),
    within_slo = countIf(duration < 500000000),  // < 500ms (em ns)
    p99_ms     = percentile(duration, 99) / 1000000,
    by: {dt.entity.service, span.name}
| fields
    service    = entityAttr(dt.entity.service, "entity.name"),
    endpoint   = span.name, requests,
    slo_pct    = round(within_slo / requests * 100, 2),
    p99_ms,
    slo_ok     = within_slo / requests >= 0.999
| filter not(slo_ok)
| sort slo_pct""",
    t=1.12, h=5.3)

# Slide 35: Caso 3
s = slide_content(prs, "Caso 3 — BizEvents (Negócio)")
code_block(s, """\
// Receita por hora — últimos 7 dias
fetch bizevents, from: now() - 7d
| filter event.type == "order.confirmed" and isNotNull(order.total)
| makeTimeseries { revenue = sum(order.total), orders = count() }, interval: 1h

// Abandono de carrinho por categoria — últimas 24h
fetch bizevents, from: now() - 24h
| filter event.type in ["cart.add","checkout.initiated"]
| summarize adds = countIf(event.type == "cart.add"),
    checkouts = countIf(event.type == "checkout.initiated"),
    by: {product.category}
| fields category = product.category,
    abandonment = round((adds - checkouts) / adds * 100, 1)
| sort abandonment, direction: "descending" """,
    t=1.12, h=5.2)

# Slide 36: Caso 4
s = slide_content(prs, "Caso 4 — Correlação CPU × Erros")
code_block(s, """\
// Hosts com CPU alta E muitos erros ao mesmo tempo
timeseries avg_cpu = avg(dt.host.cpu.usage),
    by: {dt.entity.host}, interval: 30m, from: now() - 2h
| filter avg_cpu > 70
| fields host_id = dt.entity.host, avg_cpu

| join [
    fetch logs, from: now() - 2h
    | filter loglevel == "ERROR"
    | summarize errors = count(), by: {host_id = dt.entity.host}
  ],
  on: {host_id}, kind: left

| fields host = entityAttr(host_id, "entity.name"),
    avg_cpu, errors,
    likely_cause = if(errors > 100 and avg_cpu > 85, "CPU_CORRELATED", "INVESTIGATE")
| sort avg_cpu, direction: "descending" """,
    t=1.12, h=5.3)

# Slide 37: Caso 5
s = slide_content(prs, "Caso 5 — Segurança (Brute Force)")
code_block(s, """\
// IPs com muitas tentativas de login falho — última hora
fetch logs, from: now() - 1h
| filter matchesPhrase(content, "login failed")
    or matchesPhrase(content, "authentication failed")
| parse content, "LD 'ip=' IPADDR:src_ip LD"
| filter isNotNull(src_ip)
| summarize
    attempts       = count(),
    distinct_users = countDistinct(user.id),
    by: {src_ip}
| filter attempts > 20
| sort attempts, direction: "descending" """,
    t=1.12, h=4.8)

# Slide 38: Boas Práticas
s = slide_content(prs, "Receituário de Boas Práticas")
table(s,
    headers=["Situação", "Recomendação"],
    rows=[
        ["Query exploratória (janela grande)", "scanLimitGBytes: 1  +  samplingRatio: 0.1"],
        ["Dashboard executivo", "Fixe from: e to: explícitos"],
        ["Alerta / notificação", "NUNCA use samplingRatio"],
        ["Query lenta", "Mova filter para antes do summarize"],
        ["Resolver 1–2 atributos de entidade", "Use entityAttr"],
        ["Resolver múltiplos atributos", "Use lookup"],
        ["Combinar linhas de fontes diferentes", "append (UNION ALL)"],
        ["Cruzar duas métricas / contagens", "join"],
    ],
    l=0.25, t=1.12, w=12.83, col_widths=[6.0, 6.83])


# ─── PRÓXIMOS PASSOS ─────────────────────────────────────────────────────────

# Slide 39
slide_module(prs, "", "Próximos Passos e Material")

# Slide 40: Material
s = slide_content(prs, "Material Entregue")
table(s,
    headers=["Material", "O que é"],
    rows=[
        ["Módulos (00–04)", "Referência completa de cada comando com exemplos"],
        ["Lab 01 — Fundamentos", "Exercícios de fetch, filter, timeseries"],
        ["Lab 02 — Agregação", "Exercícios de summarize, parse, makeTimeseries"],
        ["Lab 03 — Avançado", "Exercícios de lookup, join, entityAttr"],
        ["Gabaritos (3 arquivos)", "Soluções explicadas com alternativas"],
    ],
    l=0.25, t=1.12, w=12.83, col_widths=[3.8, 9.03])
note(s, "Execute os labs no Dynatrace Notebook do seu tenant.", t=5.5, size=15)

# Slide 41: Para continuar
s = slide_content(prs, "Para Continuar Aprendendo")
links = [
    "📖  Documentação Oficial DQL",
    "    docs.dynatrace.com/docs/shortlink/dql-reference",
    "",
    "🚀  DQL 301 — Advanced Concepts (WWSE)",
    "    dynatrace-wwse.github.io/enablement-dql-301/",
    "",
    "🎓  Dynatrace University  —  university.dynatrace.com",
    "",
    "💬  Community Forum  —  community.dynatrace.com",
    "",
    "Próximos temas sugeridos:",
    "   • Workflow Automations com DQL",
    "   • Davis Copilot e DQL",
    "   • Criação de Alertas e SLOs baseados em DQL",
    "   • DQL em APIs (Grail V2)",
]
tb = slide.shapes.add_textbox(Inches(0.25), Inches(1.12), Inches(12.83), Inches(6.0)) \
    if False else None  # placeholder

tb = s.shapes.add_textbox(Inches(0.25), Inches(1.12), Inches(12.83), Inches(5.9))
tf = tb.text_frame
tf.word_wrap = True
for i, line in enumerate(links):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    run = p.add_run()
    run.text = line
    run.font.name = "Calibri"
    run.font.size = Pt(15)
    run.font.color.rgb = DARK_TEXT

# Slide 42: Obrigado
slide_cover(prs,
    "Obrigado!",
    "Dúvidas e próximos passos?",
    "Alan Delcaro  ·  Dynatrace Partner Enablement")


# ─── Salvar ──────────────────────────────────────────────────────────────────

BASE = "/Users/alan.delcaro/Library/CloudStorage/OneDrive-Dynatrace/Documents/Profissional/Testes_e_Dev/workshop_parceiros"
output = f"{BASE}/slides/apresentacao-dql-editavel.pptx"
prs.save(output)
print(f"✅  PPTX editável salvo em: {output}")
print(f"    Total de slides: {len(prs.slides)}")
